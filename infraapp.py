import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from prophet import Prophet
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from dotenv import load_dotenv
from groq import Groq
import os

# Load environment variables
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

if not GROQ_API_KEY:
    st.error("🚨 API Key is missing! Set it in Streamlit Secrets or a .env file.")
    st.stop()

# Streamlit page setup
st.set_page_config(page_title="GCP Cost & Usage Forecasting", page_icon="📊", layout="wide")
st.title("📊 GCP Forecasting Agent")
st.markdown("Upload your Excel file to forecast `Total Cost` and `Total Usage` by `Product` group.")

# File uploader
uploaded_file = st.file_uploader("📤 Upload Excel File", type=["xlsx"])
forecast_period = st.slider("📅 Forecast Horizon (months)", 1, 36, 12)

if uploaded_file:
    df = pd.read_excel(uploaded_file)

    # Clean & prepare data
    df['Period'] = pd.to_datetime(df['Period'])
    df['Total Cost'] = pd.to_numeric(df['Total Cost'], errors='coerce')
    df['Total Usage'] = pd.to_numeric(df['Total Usage'], errors='coerce')
    df = df.dropna(subset=['Period', 'Product', 'Total Cost', 'Total Usage'])

    products = df['Product'].dropna().unique()
    client = Groq(api_key=GROQ_API_KEY)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = "GCP Forecast Report"

    for metric in ['Total Cost', 'Total Usage']:
        for product in products:
            st.subheader(f"📈 {metric} Forecast – {product}")
            df_product = df[df['Product'] == product][['Period', metric]].rename(columns={'Period': 'ds', metric: 'y'})

            if df_product.empty or df_product['y'].sum() == 0:
                st.warning(f"No valid {metric} data for {product}")
                continue

            model = Prophet()
            model.fit(df_product)
            future = model.make_future_dataframe(periods=forecast_period, freq='M')
            forecast = model.predict(future)

            fig = model.plot(forecast)
            st.pyplot(fig)

            # Save figure to PPT
            buf = BytesIO()
            fig.savefig(buf, format='png')
            buf.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"{metric} Forecast – {product}"
            slide.shapes.add_picture(buf, Inches(1), Inches(1.2), width=Inches(8.5))

            # AI Commentary
            ai_data = forecast[['ds', 'yhat', 'yhat_lower', 'yhat_upper']].tail(24).to_json(orient='records', date_format='iso')
            prompt = f"""
            You are an FP&A leader at a tech company analyzing GCP {metric} for {product}.
            Based on this JSON forecast data, summarize:
            - Key trends and seasonality
            - Risks and recommendations
            - A CFO-style summary using the Pyramid Principle

            Data: {ai_data}
            """

            response = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a SaaS FP&A expert."},
                    {"role": "user", "content": prompt}
                ],
                model="llama3-8b-8192"
            )
            commentary = response.choices[0].message.content
            st.markdown(f"### 🧠 AI Commentary – {product} ({metric})")
            st.write(commentary)

            # Add commentary to PPT
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{product} – {metric} Insights"
            slide.placeholders[1].text = commentary

    # Export PPT
    ppt_buf = BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)

    st.download_button("📥 Download Executive PowerPoint", ppt_buf, "GCP_Forecast_Summary.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")

else:
    st.info("Please upload an Excel file with 'Period', 'Product', 'Total Cost', and 'Total Usage'.")
